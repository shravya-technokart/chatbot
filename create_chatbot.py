import os
import fitz 
import pymongo
import docx2txt
from PIL import Image  
from openai import OpenAI
from pptx import Presentation
from pytesseract import pytesseract

from dotenv import load_dotenv
load_dotenv()

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY")
)
mongo_url = "mongodb+srv://shravya:shetty34@chatbot-details.icmkkzh.mongodb.net"
client_mongo = pymongo.MongoClient(mongo_url)
db = client_mongo['chatbot-details']
collection = db['docs-data']
collection_chatbot = db['chatbot-details']
collection_manager = db['manager-details']

# chat_bot = "Trainer"
manager_email = "manager@gmail.com"

class ChatBot:
    def __init__(self):
        self.chat_bot = input("Enter Chatbot Name: ")
        self.manager_email = "manager@gmail.com"
        collection_chatbot.insert_one(
            {
                'manager_email': self.manager_email,
                'chat_bot': self.chat_bot,
                'system_prompt': "You are a question answering chatbot."
            }
        )
        self.personality = ""
        self.company = ""
        self.services = ""
        self.others = ""
        self.instructions = ""

    def get_domain_knowledge(self) -> dict:
        personality = input("Personality of Chatbot: ")
        if personality != "":
            self.personality = f"""
                As a chatbot you need to maintain your PERSONALITY as described below.
                PERSONALITY - '{personality}'
            ----------------------------------------------------------------------------------------------------
            """
        company = input("Company Details: ")
        if company!="":
            self.company = f"""
                Company_Details under which you are acting as a chatbot is described below.
                Company_Details - '{company}'
            ----------------------------------------------------------------------------------------------------
            """
        services = input("Service Details: ")
        if services!="":
            self.services = f"""
                Service_Details which your company provides is described below.
                Service_Details - '{services}'
            ----------------------------------------------------------------------------------------------------
            """
        other = input("Other Details : ")
        if other!="":
            self.other = f"""
                Other_Details which your company provides is described below.
                Other_Details - '{other}'
        ----------------------------------------------------------------------------------------------------
            """
        instructions = input("Instructions : ")
        if instructions!="":
            self.instructions = f"""
                Instructions are provided below which you should follow to respond user -
                Intructions - {instructions}
        """
        # return {"personality":self.personality, "company":self.company, "services":self.services, "other_details": self.others}

    def get_documents(self):
        def image_to_text(img_path):
            img = Image.open(img_path) 
            img_text = pytesseract.image_to_string(img) 
            return img_text

        def pdf_to_text(pdf_path):
            pdf_document = fitz.open(pdf_path)
            pdf_text = ""
            for page in pdf_document:
                pdf_text+=page.get_text()
            return pdf_text 

        def txt_to_text(txt_path):
            with open(txt_path, 'r') as f:
                text = f.read()
            return text 

        def word_to_text(word_path):
            word_text = docx2txt.process(word_path)
            return word_text 

        def ppt_to_text(ppt_path):
            prs = Presentation(ppt_path)
            ppt_data = ""
            for slide in prs.slides:
                list_of_elements = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        list = ""
                        for run in paragraph.runs:
                            list += run.text
                        list_of_elements.append(list)
                for elements in list_of_elements:
                    ppt_data += elements 
            return ppt_data

        for file_path in os.listdir('./pdfs'):
            file_name, file_extension = os.path.splitext(file_path)
            image_extension = [".jpg", ".jpeg", ".png", ".tiff", ".gif", ".bmp"]
            path = os.path.join('./pdfs', file_path)
            if file_extension==".pdf":
                sales_data = pdf_to_text(path)
            elif file_extension in image_extension:
                sales_data = image_to_text(path)
            elif file_extension==".txt":
                sales_data = txt_to_text(path)
            elif file_extension==".docx":
                sales_data = word_to_text(path)
            elif file_extension==".pptx":
                sales_data = ppt_to_text(path)
            else:
                print('Improper document format')
                sales_data = ""
        
            def store_chunk(chunk):
                def get_embedding(text, model="text-embedding-3-small"):
                    text = text.replace("\n", " ")
                    return client.embeddings.create(input=text, model=model, dimensions=1536).data[0].embedding

                embedded_chunk = get_embedding(chunk, model='text-embedding-3-small')

                collection.insert_one(
                    {
                    'chat_bot': self.chat_bot,
                    'manager_email': manager_email,
                    'chunk': chunk,
                    'embedded_chunk': embedded_chunk
                    }
                )

            i = 0
            sales_data_words = sales_data.split(' ')
            while i<len(sales_data_words):
                chunk_words = sales_data_words[i:i+500]
                i += 250
                chunk = ' '.join(chunk_words)
                store_chunk(chunk)       

    def generate_system_message(self) -> str:
        system_prompt = f"""
            You are a chatbot designed to reply user based on Context that would be given to you. 
            
            {self.instructions}
            
            Some of the Characteristic Information for you is given below - 
            {self.personality}
            {self.company}
            {self.services}
            {self.others}
        """
        collection_chatbot.update_one({ "manager_email": manager_email, "chat_bot": self.chat_bot }, { "$set": { "system_prompt": system_prompt } })
        
        print('System Prompt added succesfully.')
    
chatbot_manager = ChatBot()
chatbot_manager.get_documents()
chatbot_manager.get_domain_knowledge()
chatbot_manager.generate_system_message()