import os
import fitz 
import pymongo
import docx2txt
from PIL import Image  
from openai import OpenAI
from pptx import Presentation
from pytesseract import pytesseract

from dotenv import load_dotenv
load_dotenv()

client = OpenAI(
    api_key=os.environ.get("OPENAI_API_KEY"),
)
client_mongo = pymongo.MongoClient(os.environ.get("mongodb_url"))

db = client_mongo['chatbot-details']
collection = db['docs-data']

# chat_bot = "Trainer"
# manager_email = "manager@gmail.com"

def get_documents(chat_bot, manager_email):
    def image_to_text(img_path):
        img = Image.open(img_path) 
        img_text = pytesseract.image_to_string(img) 
        return img_text

    def pdf_to_text(pdf_path):
        pdf_document = fitz.open(pdf_path)
        pdf_text = ""
        for page in pdf_document:
            pdf_text+=page.get_text()
        return pdf_text 

    def txt_to_text(txt_path):
        with open(txt_path, 'r') as f:
            text = f.read()
        return text 

    def word_to_text(word_path):
        word_text = docx2txt.process(word_path)
        return word_text 

    def ppt_to_text(ppt_path):
        prs = Presentation(ppt_path)
        ppt_data = ""
        for slide in prs.slides:
            list_of_elements = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    list = ""
                    for run in paragraph.runs:
                        list += run.text
                    list_of_elements.append(list)
            for elements in list_of_elements:
                ppt_data += elements 
        return ppt_data

    for file_path in os.listdir('./pdfs'):
        file_name, file_extension = os.path.splitext(file_path)
        image_extension = [".jpg", ".jpeg", ".png", ".tiff", ".gif", ".bmp"]
        path = os.path.join('./pdfs', file_path)
        if file_extension==".pdf":
            sales_data = pdf_to_text(path)
        elif file_extension in image_extension:
            sales_data = image_to_text(path)
        elif file_extension==".txt":
            sales_data = txt_to_text(path)
        elif file_extension==".docx":
            sales_data = word_to_text(path)
        elif file_extension==".pptx":
            sales_data = ppt_to_text(path)
        else:
            print('Improper document format')
            sales_data = ""
        
        def store_chunk(chunk):
            def get_embedding(text, model="text-embedding-3-small"):
                text = text.replace("\n", " ")
                return client.embeddings.create(input=text, model=model, dimensions=1536).data[0].embedding

            embedded_chunk = get_embedding(chunk, model='text-embedding-3-small')

            collection.insert_one(
                {
                    'chat_bot': chat_bot,
                    'manager_email': manager_email,
                    'chunk': chunk,
                    'embedded_chunk': embedded_chunk
                }
            )

        i = 0
        sales_data_words = sales_data.split(' ')
        while i<len(sales_data_words):
            chunk_words = sales_data_words[i:i+1000]
            i += 500
            chunk = ' '.join(chunk_words)
            store_chunk(chunk)